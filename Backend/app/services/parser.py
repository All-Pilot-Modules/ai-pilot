from pptx import Presentation
import pdfplumber
import docx
import os

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pptx":
        prs = Presentation(file_path)
        return "\n".join([
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        ])

    elif ext == ".docx":
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

    elif ext == ".pdf":
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        print(text)
        return text

    else:
        raise ValueError("Unsupported file type for text extraction.")