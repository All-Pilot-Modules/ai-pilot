"""
Unified text extractor for multiple file formats
Supports: PDF, DOCX, PPTX, TXT
Uses LlamaParse for testbank extraction (AI-powered)
"""
import os
from typing import Dict, Any
import fitz  # PyMuPDF for PDF
from docx import Document as DocxDocument  # python-docx for DOCX
from pptx import Presentation  # python-pptx for PPTX
from llama_parse import LlamaParse


def extract_text_from_pdf(file_path: str) -> Dict[str, Any]:
    """
    Extract text from PDF file

    Args:
        file_path: Path to PDF file

    Returns:
        {
            'text': str,  # Full text
            'metadata': {
                'pages': int,
                'page_texts': list[str]  # Text per page
            }
        }
    """
    doc = fitz.open(file_path)
    full_text = ""
    page_texts = []

    for page_num, page in enumerate(doc, start=1):
        page_text = page.get_text()
        page_texts.append(page_text)
        full_text += f"\n--- Page {page_num} ---\n{page_text}"

    doc.close()

    return {
        'text': full_text.strip(),
        'metadata': {
            'pages': len(page_texts),
            'page_texts': page_texts
        }
    }


def extract_text_from_docx(file_path: str) -> Dict[str, Any]:
    """
    Extract text from DOCX file

    Args:
        file_path: Path to DOCX file

    Returns:
        {
            'text': str,
            'metadata': {
                'paragraphs': int
            }
        }
    """
    doc = DocxDocument(file_path)
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
    full_text = "\n\n".join(paragraphs)

    return {
        'text': full_text.strip(),
        'metadata': {
            'paragraphs': len(paragraphs)
        }
    }


def extract_text_from_pptx(file_path: str) -> Dict[str, Any]:
    """
    Extract text from PPTX file

    Args:
        file_path: Path to PPTX file

    Returns:
        {
            'text': str,
            'metadata': {
                'slides': int,
                'slide_texts': list[str]
            }
        }
    """
    prs = Presentation(file_path)
    full_text = ""
    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text += shape.text + "\n"

        slide_texts.append(slide_text.strip())
        full_text += f"\n--- Slide {slide_num} ---\n{slide_text}"

    return {
        'text': full_text.strip(),
        'metadata': {
            'slides': len(slide_texts),
            'slide_texts': slide_texts
        }
    }


def extract_text_from_txt(file_path: str) -> Dict[str, Any]:
    """
    Extract text from TXT file

    Args:
        file_path: Path to TXT file

    Returns:
        {
            'text': str,
            'metadata': {
                'lines': int
            }
        }
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()

    return {
        'text': text.strip(),
        'metadata': {
            'lines': len(text.splitlines())
        }
    }


def extract_text_with_llamaparse(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    Extract text using LlamaParse AI-powered extraction
    Handles complex layouts, tables, multi-column formats, and scanned PDFs

    Args:
        file_path: Path to PDF or DOCX file
        file_type: File extension (pdf or docx)

    Returns:
        {
            'text': str,
            'metadata': {
                'pages': int,
                'extraction_method': 'llamaparse'
            }
        }

    Raises:
        Exception: If LlamaParse extraction fails
    """
    try:
        # Get API key from environment
        api_key = os.getenv('LLAMA_CLOUD_API_KEY')
        if not api_key:
            raise ValueError("LLAMA_CLOUD_API_KEY not found in environment variables")

        # Initialize LlamaParse
        parser = LlamaParse(
            api_key=api_key,
            result_type="text",  # Return plain text
            verbose=False
        )

        # Parse the document
        documents = parser.load_data(file_path)

        # Combine all pages/sections into single text
        full_text = "\n\n".join([doc.text for doc in documents])

        return {
            'text': full_text.strip(),
            'metadata': {
                'pages': len(documents),
                'extraction_method': 'llamaparse'
            }
        }

    except Exception as e:
        print(f"LlamaParse extraction failed: {str(e)}")
        raise


def extract_text_from_file(file_path: str, file_type: str, is_testbank: bool = False) -> Dict[str, Any]:
    """
    Unified text extractor - automatically detects file type

    Args:
        file_path: Path to file
        file_type: File extension (pdf, docx, pptx, txt)
        is_testbank: If True, uses LlamaParse for PDF/DOCX extraction

    Returns:
        {
            'text': str,
            'metadata': dict
        }

    Raises:
        ValueError: If file type is not supported
    """
    file_type = file_type.lower()

    # Use LlamaParse for testbank PDFs and DOCX files
    if is_testbank and file_type in ['pdf', 'docx', 'doc']:
        try:
            print(f"Using LlamaParse for testbank extraction: {file_path}")
            return extract_text_with_llamaparse(file_path, file_type)
        except Exception as e:
            print(f"LlamaParse failed, falling back to standard extractor: {str(e)}")
            # Fall back to standard extractors if LlamaParse fails

    # Standard extractors
    if file_type == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_type in ['docx', 'doc']:
        return extract_text_from_docx(file_path)
    elif file_type in ['pptx', 'ppt']:
        return extract_text_from_pptx(file_path)
    elif file_type == 'txt':
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}. Supported: pdf, docx, pptx, txt")
